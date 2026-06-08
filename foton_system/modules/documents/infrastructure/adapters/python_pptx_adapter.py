import os
import re
from typing import Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.table import Table, _Cell, _Row
from foton_system.modules.shared.infrastructure.config.logger import setup_logger
from foton_system.modules.documents.application.ports.document_service_port import DocumentServicePort

logger = setup_logger()


class PythonPPTXAdapter(DocumentServicePort):
    """Adapter para manipulação de apresentações PPTX via python-pptx."""

    def __init__(self) -> None:
        pass

    def load_document(self, path: str) -> Presentation:
        """Carrega um arquivo PPTX do disco."""
        if not os.path.exists(path):
            raise FileNotFoundError(f"Arquivo PPTX não encontrado: {path}")
        return Presentation(path)

    def save_document(self, document: Presentation, path: str) -> None:
        """Persiste a apresentação PPTX no disco."""
        try:
            document.save(path)
            logger.info(f"Apresentação salva em: {path}")
        except Exception as e:
            logger.error(f"Erro ao salvar apresentação: {e}")
            raise

    def replace_text(self, document: Presentation, replacements: Dict[str, str]) -> Presentation:
        """Substitui chaves por valores em todos os slides e tabelas."""
        logger.info('Iniciando substituição de textos no PPTX...')

        for slide in document.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    self._replace_in_text_frame(shape.text_frame, replacements)
                if shape.has_table:
                    self._replace_in_table(shape.table, replacements)

        logger.info('Substituição de textos concluída.')
        return document

    def _replace_in_text_frame(self, text_frame: Any, replacements: Dict[str, str]) -> None:
        """Substitui chaves no text frame, consolidando runs antes."""
        for paragraph in text_frame.paragraphs:
            self._consolidate_runs(paragraph)
            for run in paragraph.runs:
                run.text = self._replace_keys_in_text(run.text, replacements)

    def _consolidate_runs(self, paragraph: Any) -> None:
        """Mescla runs do parágrafo para evitar chaves quebradas."""
        text: str = "".join(run.text for run in paragraph.runs)
        if '@' in text and len(paragraph.runs) > 1:
            paragraph.runs[0].text = text
            for run in paragraph.runs[1:]:
                run.text = ""

    def _replace_in_table(self, table: Table, replacements: Dict[str, str]) -> None:
        """Substitui chaves em todas as células da tabela."""
        for row in table.rows:
            for cell in row.cells:
                self._replace_in_text_frame(cell.text_frame, replacements)

    def _replace_keys_in_text(self, text: str, replacements: Dict[str, str]) -> str:
        """Substitui chaves no texto usando regex case-insensitive."""
        sorted_keys = sorted(replacements.keys(), key=len, reverse=True)
        for key in sorted_keys:
            pattern = r'(?<![\w.])' + re.escape(key) + r'(?!\.[a-z]{2,}\b)'
            new_val: str = str(replacements[key])
            text = re.sub(pattern, new_val, text, flags=re.IGNORECASE)
        return text
