import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from foton_system.core.logger import setup_logger

logger = setup_logger()

class PPTXHandler:
    def __init__(self):
        pass

    def load_presentation(self, path):
        if not os.path.exists(path):
            raise FileNotFoundError(f"Arquivo PPTX não encontrado: {path}")
        return Presentation(path)

    def save_presentation(self, presentation, path):
        try:
            presentation.save(path)
            logger.info(f"Apresentação salva em: {path}")
        except Exception as e:
            logger.error(f"Erro ao salvar apresentação: {e}")
            raise

    def replace_text(self, presentation, replacements):
        """
        Replaces keys with values in a PowerPoint presentation.
        """
        logger.info('Iniciando substituição de textos no PPTX...')
        
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    self._replace_in_text_frame(shape.text_frame, replacements)

                if shape.has_table:
                    self._replace_in_table(shape.table, replacements)
        
        logger.info('Substituição de textos concluída.')
        return presentation

    def _replace_in_text_frame(self, text_frame, replacements):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = self._replace_keys_in_text(run.text, replacements)

    def _replace_in_table(self, table, replacements):
        for row in table.rows:
            for cell in row.cells:
                self._replace_in_text_frame(cell.text_frame, replacements)

    def _replace_keys_in_text(self, text, replacements):
        import re
        # Sort keys by length (descending) to avoid partial replacement issues
        sorted_keys = sorted(replacements.keys(), key=len, reverse=True)
        
        for key in sorted_keys:
            if key in text:
                # Use regex to ensure we don't replace inside words/emails
                # Lookbehind: (?<![\w.]) checks if previous char is NOT alphanumeric/underscore/dot
                # Lookahead: (?!\.[a-z]{2,}\b) checks if NOT followed by domain-like suffix
                pattern = r'(?<![\w.])' + re.escape(key) + r'(?!\.[a-z]{2,}\b)'
                new_val = str(replacements[key])
                text = re.sub(pattern, new_val, text)
                
        return text
