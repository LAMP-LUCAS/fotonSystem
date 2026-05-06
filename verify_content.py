import sys
from docx import Document
from pptx import Presentation

def verify_docx(path):
    print(f"\n--- Verificando DOCX: {path} ---")
    doc = Document(path)
    for p in doc.paragraphs:
        if "SIMONE" in p.text or "PROJETO" in p.text or "9.347,65" in p.text:
            print(f"Encontrado: {p.text}")

def verify_pptx(path):
    print(f"\n--- Verificando PPTX: {path} ---")
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                if "SIMONE" in text or "PROJETO" in text or "9.347,65" in text:
                    print(f"Encontrado: {text}")

if __name__ == "__main__":
    verify_docx(r'C:\Users\Lucas\OneDrive\LAMP_ARQUITETURA\CLIENTES\SIMONE_SEBASTIAO_TESTE\GERADO_02-COD_DOC_CT_00_R00_PROPOSTA-PROJETO.docx')
    verify_pptx(r'C:\Users\Lucas\OneDrive\LAMP_ARQUITETURA\CLIENTES\SIMONE_SEBASTIAO_TESTE\GERADO_02-COD_DOC_PC_00_R00_PROPOSTA_GENERICO.pptx')
